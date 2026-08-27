"""Skill tools — the ``read_skill`` tool used by the model to load skill content.

This module provides a factory :func:`make_skill_tools` that returns a
``read_skill`` tool closed over a skills directory path.  The tool reads
files directly from the host filesystem (skills are read-only reference
content), so it works identically for both local-subprocess and
Docker-container sandbox providers — no sandbox mount required.

Binary file support
-------------------
The tool can extract text content from ``.docx`` and ``.pptx`` files
so that LLMs can consume reference documents directly.  The required
libraries (``python-docx``, ``python-pptx``) are soft dependencies —
when they are not installed the tool returns a clear error message
telling the caller how to install the missing package.
"""

from __future__ import annotations

import logging
import re
from collections.abc import Awaitable, Callable
from pathlib import Path
from typing import Any

from langchain.tools import ToolRuntime, tool

from agent_sdk.sandbox.base import SandboxProvider
from agent_sdk.skills.loader import load_skills
from agent_sdk.skills.types import Skill

logger = logging.getLogger(__name__)

# Regex matching YAML frontmatter blocks (---\n...\n---) in SKILL.md files.
# Used to strip metadata before returning content to the agent — the
# name/description/license fields are for the loader, not the model.
_FRONTMATTER_RE = re.compile(r"^---\s*\r?\n.*?\r?\n---\s*\r?\n", re.DOTALL)

# File extensions that indicate a skill contains binary (non-text) resources.
# When a skill includes ANY file with one of these extensions, the entire
# skill directory is injected into the sandbox on first read_skill access.
_BINARY_EXTENSIONS: frozenset[str] = frozenset(
    {
        # Office documents
        ".docx",
        ".pptx",
        ".xlsx",
        ".pdf",
        ".doc",
        ".ppt",
        ".xls",
        # Images
        ".png",
        ".jpg",
        ".jpeg",
        ".gif",
        ".bmp",
        ".svg",
        ".ico",
        ".webp",
        # Archives
        ".zip",
        ".tar",
        ".gz",
        ".bz2",
        ".xz",
        ".7z",
        # Fonts
        ".ttf",
        ".otf",
        ".woff",
        ".woff2",
        # Compiled / native
        ".pyc",
        ".pyd",
        ".so",
        ".dll",
        ".dylib",
        ".exe",
        # Database
        ".db",
        ".sqlite",
        ".sqlite3",
        # Other binary
        ".bin",
        ".dat",
    }
)

# Per-skill cache: skill_name → has_binary_files (populated lazily, cleared
# when skills directory changes — i.e. never at runtime since skills are
# read-only reference content).
_skill_has_binary_cache: dict[str, bool] = {}

# Runtime type alias — matches the sandbox tools convention.
_RuntimeType = ToolRuntime  # type: ignore[assignment]

#: Async callback injected by the app layer to decide whether a skill is in
#: the current user's available set (built-in skills are resolved from the
#: filesystem *before* this callback is consulted, so it only ever sees
#: custom/OBS-resident skill names).
IsAvailableCallback = Callable[[str], Awaitable[bool]]

#: Async callback injected by the app layer to download every file of a
#: custom skill from object storage.  Returns ``(rel_path, bytes)`` pairs,
#: where ``rel_path`` is relative to the skill root (e.g. ``SKILL.md``,
#: ``scripts/run.py``).
FetchSkillFilesCallback = Callable[[str], Awaitable[list[tuple[str, bytes]]]]

#: Async callback injected by the app layer that returns the current user's
#: personal skills as ``(name, description)`` pairs (own skills + added &
#: approved marketplace skills).
ListPersonalSkillsCallback = Callable[[], Awaitable[list[tuple[str, str]]]]


def _get_sandbox_from_runtime(
    runtime: ToolRuntime | None,
    sandbox_provider: SandboxProvider | None,
) -> Any | None:
    """Extract the sandbox instance from *runtime* state, or ``None``.

    Mirrors the logic from :func:`agent_sdk.sandbox.tools._try_get_sandbox`
    but is self-contained to avoid a circular dependency.
    """
    if runtime is None or runtime.state is None or sandbox_provider is None:
        return None
    sandbox_state = runtime.state.get("sandbox")
    if sandbox_state is None:
        return None
    sandbox_id = sandbox_state.get("sandbox_id")
    if sandbox_id is None:
        return None
    return sandbox_provider.get(sandbox_id)


def _skill_has_binary_files(skill_dir: Path) -> bool:
    """Check whether *skill_dir* contains any binary (non-text) files.

    Results are cached per skill name (skills are read-only at runtime,
    so the answer never changes within a process lifetime).
    """
    name = skill_dir.name
    if name in _skill_has_binary_cache:
        return _skill_has_binary_cache[name]

    for f in skill_dir.rglob("*"):
        if f.is_file() and f.suffix.lower() in _BINARY_EXTENSIONS:
            _skill_has_binary_cache[name] = True
            return True

    _skill_has_binary_cache[name] = False
    return False


def _ensure_sandbox_for_injection(
    runtime: ToolRuntime,
    sandbox_provider: SandboxProvider,
) -> Any | None:
    """Get or acquire a sandbox for skill file injection.

    Mirrors the lazy-acquire logic from
    :func:`agent_sdk.sandbox.tools._ensure_sandbox` but is
    self-contained to avoid a circular dependency.  Returns the
    sandbox instance or ``None`` if a thread_id cannot be resolved
    (e.g. the conversation hasn't started yet).
    """
    # 1) Already bound in runtime state — just return it.
    sandbox = _get_sandbox_from_runtime(runtime, sandbox_provider)
    if sandbox is not None:
        return sandbox

    # 2) Not yet bound — need to acquire a sandbox for this thread.
    #    The previous sandbox container may have been destroyed (warm
    #    pool idle timeout → Pod deleted → emptyDir gone).  Clear the
    #    per-conversation skill-injection markers so that binary skill
    #    files are re-injected into the fresh container.
    if runtime.state is not None:
        runtime.state.pop("injected_skills", None)
    thread_id: str | None = None
    if runtime.context is not None:
        thread_id = runtime.context.get("thread_id")
    if thread_id is None and runtime.config is not None:
        thread_id = runtime.config.get("configurable", {}).get("thread_id")
    if thread_id is None:
        # Last resort: extract from thread_data.workspace_path.
        thread_data = runtime.state.get("thread_data") if runtime.state is not None else None
        if thread_data is not None:
            from agent_sdk.utils.thread import extract_thread_id

            thread_id = extract_thread_id(thread_data)

    if thread_id is None:
        return None

    try:
        sandbox_id = sandbox_provider.acquire(thread_id)
    except Exception:
        logger.warning("Failed to acquire sandbox for skill injection", exc_info=True)
        return None

    sandbox = sandbox_provider.get(sandbox_id)
    if sandbox is not None and runtime.state is not None:
        runtime.state["sandbox"] = {"sandbox_id": sandbox_id}
    return sandbox


def _inject_skill_files(
    skill_dir: Path,
    sandbox_provider: SandboxProvider,
    runtime: ToolRuntime,
) -> str:
    """Copy all files from *skill_dir* into the sandbox.

    Only called when the skill contains binary files and hasn't been
    injected yet for this conversation (tracked via ``runtime.state``).

    The destination path depends on the sandbox type:

    * **LocalSandbox** — relative ``.skills/<name>/``, resolved
      against ``LocalSandbox._workspace`` (the host workspace root).
    * **AioSandbox (Docker)** — absolute
      ``/mnt/user-data/workspace/.skills/<name>/`` so that the agent
      sees the files under its regular workspace.

    Returns a human-readable injection summary.
    """
    sandbox = _ensure_sandbox_for_injection(runtime, sandbox_provider)
    if sandbox is None:
        return "\n\n⚠️  此 skill 包含二进制文件，但无法获取 sandbox，文件未注入。 请先执行一个 sandbox 操作（如 ls）初始化环境后重试。"

    from agent_sdk.sandbox.local.provider import LocalSandbox

    skill_name = skill_dir.name
    is_local = isinstance(sandbox, LocalSandbox)
    skill_path_hint = f".skills/{skill_name}/" if is_local else f"/mnt/user-data/workspace/.skills/{skill_name}/"

    logger.info(
        "Injecting skill '%s' (%s files) into sandbox at %s",
        skill_name,
        "local" if is_local else "docker",
        skill_path_hint,
    )

    injected_count = 0

    for f in skill_dir.rglob("*"):
        if f.is_file() and not f.name.startswith("."):
            rel_path = f.relative_to(skill_dir).as_posix()
            if is_local:
                dest = f".skills/{skill_name}/{rel_path}"
            else:
                dest = f"/mnt/user-data/workspace/.skills/{skill_name}/{rel_path}"
            try:
                content_bytes = f.read_bytes()
                sandbox.update_file(dest, content_bytes)
                injected_count += 1
                logger.debug("  injected: %s → %s", rel_path, dest)
            except Exception:
                logger.warning(
                    "Failed to inject skill file %s into sandbox",
                    f,
                    exc_info=True,
                )

    # Mark as injected so we don't re-inject on subsequent calls.
    if runtime.state is not None:
        injected: set[str] = runtime.state.get("injected_skills")
        if injected is None:
            injected = set()
            runtime.state["injected_skills"] = injected
        injected.add(skill_name)

    logger.info(
        "Skill '%s' injected: %d files → %s",
        skill_name,
        injected_count,
        skill_path_hint,
    )

    return f"\n\n📦 此 skill 包含二进制文件，已将全部 {injected_count} 个文件注入 sandbox: `{skill_path_hint}`\n可直接用 bash 执行脚本或引用模板文件。"


def _inject_skill_file_list(
    files: list[tuple[str, bytes]],
    skill_name: str,
    sandbox_provider: SandboxProvider,
    runtime: ToolRuntime,
) -> str:
    """Write in-memory skill files into the agent's sandbox workspace.

    Used for custom skills fetched from object storage (OBS): the agent's
    filesystem tools cannot see OBS, so the skill's files must be
    materialised under ``.skills/<name>/`` before the agent can execute
    scripts or reference templates.  Injection is per-conversation
    (tracked via ``runtime.state``) and happens once per skill.
    """
    sandbox = _ensure_sandbox_for_injection(runtime, sandbox_provider)
    if sandbox is None:
        return "\n\n⚠️  此 skill 的文件无法注入 sandbox（sandbox 不可用）。 请先执行一个 sandbox 操作（如 ls）初始化环境后重试。"

    from agent_sdk.sandbox.local.provider import LocalSandbox

    is_local = isinstance(sandbox, LocalSandbox)
    skill_path_hint = f".skills/{skill_name}/" if is_local else f"/mnt/user-data/workspace/.skills/{skill_name}/"

    logger.info(
        "Injecting custom skill '%s' (%d files) into sandbox at %s",
        skill_name,
        len(files),
        skill_path_hint,
    )

    injected_count = 0
    for rel_path, content_bytes in files:
        dest = f".skills/{skill_name}/{rel_path}" if is_local else f"/mnt/user-data/workspace/.skills/{skill_name}/{rel_path}"
        try:
            sandbox.update_file(dest, content_bytes)
            injected_count += 1
            logger.debug("  injected: %s → %s", rel_path, dest)
        except Exception:
            logger.warning(
                "Failed to inject custom skill file %s into sandbox",
                rel_path,
                exc_info=True,
            )

    # Mark as injected so we don't re-inject on subsequent calls.
    if runtime.state is not None:
        injected: set[str] = runtime.state.get("injected_skills")
        if injected is None:
            injected = set()
            runtime.state["injected_skills"] = injected
        injected.add(skill_name)

    return f"\n\n📦 已将自定义技能 {skill_name} 的 {injected_count} 个文件注入 sandbox: `{skill_path_hint}`\n可直接用 bash 执行脚本或引用模板文件。"


# File extensions recognised as plain text — opened with UTF-8.
_TEXT_EXTENSIONS: frozenset[str] = frozenset(
    {
        ".md",
        ".py",
        ".js",
        ".ts",
        ".jsx",
        ".tsx",
        ".txt",
        ".yaml",
        ".yml",
        ".json",
        ".html",
        ".css",
        ".xml",
        ".csv",
        ".sh",
        ".bash",
        ".toml",
        ".ini",
        ".cfg",
        ".rst",
        ".java",
        ".go",
        ".rs",
        ".c",
        ".cpp",
        ".h",
        ".rb",
        ".php",
        ".sql",
        ".r",
        ".m",
        ".swift",
        ".kt",
        ".scala",
        ".lua",
        ".pl",
        ".pm",
    }
)

# Extensions for which we have specialised text-extraction logic.
_DOCX_EXTENSIONS: frozenset[str] = frozenset({".docx"})
_PPTX_EXTENSIONS: frozenset[str] = frozenset({".pptx"})


# =================================================================
# Binary-file readers (soft-dependency — fail gracefully on ImportError)
# =================================================================


def _read_docx(path: Path) -> str:
    """Extract text content from a ``.docx`` file.

    Returns paragraphs joined by newlines; tables are rendered as
    pipe-delimited rows.
    """
    try:
        from docx import Document  # type: ignore[import-untyped]
    except ImportError:
        return "Error: python-docx is not installed. Install it with: pip install python-docx"

    try:
        doc = Document(str(path))
    except Exception as exc:
        return f"Error: Failed to open .docx file: {exc}"

    parts: list[str] = []

    for para in doc.paragraphs:
        if para.text.strip():
            parts.append(para.text)

    for i, table in enumerate(doc.tables):
        parts.append(f"\n[表格 {i + 1}]")
        for row in table.rows:
            cells = [cell.text.replace("\n", " ") for cell in row.cells]
            parts.append(" | ".join(cells))

    return "\n".join(parts)


def _read_pptx(path: Path) -> str:
    """Extract text content from a ``.pptx`` file.

    Returns slide-by-slide output with shape text grouped per slide.
    """
    try:
        from pptx import Presentation  # type: ignore[import-untyped]
    except ImportError:
        return "Error: python-pptx is not installed. Install it with: pip install python-pptx"

    try:
        prs = Presentation(str(path))
    except Exception as exc:
        return f"Error: Failed to open .pptx file: {exc}"

    parts: list[str] = []
    for i, slide in enumerate(prs.slides):
        slide_lines: list[str] = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        slide_lines.append(text)
        if slide_lines:
            parts.append(f"\n=== Slide {i + 1} ===")
            parts.extend(slide_lines)

    return "\n".join(parts)


# =================================================================
# File-type dispatch
# =================================================================


def _read_skill_file(target_file: Path) -> str:
    """Read a file in a skill directory, dispatching on file extension.

    Plain-text files are read with UTF-8.  ``.docx`` and ``.pptx``
    files have their text content extracted.  Other binary files
    return a descriptive error.
    """
    suffix = target_file.suffix.lower()

    # ── specialised binary readers ──
    if suffix in _DOCX_EXTENSIONS:
        return _read_docx(target_file)
    if suffix in _PPTX_EXTENSIONS:
        return _read_pptx(target_file)

    # ── plain text ──
    try:
        return target_file.read_text(encoding="utf-8")
    except UnicodeDecodeError:
        return f"Error: Cannot read '{target_file.name}' as text — it appears to be a binary file. Supported binary formats: .docx, .pptx"
    except OSError as exc:
        return f"Error: Failed to read file '{target_file.name}': {exc}"


def _read_builtin_skill(
    target: Skill,
    file: str | None,
    sandbox_provider: SandboxProvider | None,
    runtime: ToolRuntime | None,
) -> str:
    """Read a built-in skill's ``SKILL.md`` or a supporting file from disk.

    Extracted from :func:`_build_read_skill_tool` — handles the whole
    filesystem read path (file resolution, directory listing, docx/pptx
    extraction, frontmatter stripping, and binary-file sandbox injection).
    """
    name = target.name

    # Resolve which file to read.
    if file is not None:
        # Normalise separators and reject traversal.
        clean = file.replace("\\", "/")
        for segment in clean.split("/"):
            if segment == "..":
                return f"Error: Path traversal is not allowed: '{file}'. Use a relative path within the skill directory (e.g. 'scripts/generate.js')."
        target_file = (target.skill_dir / clean).resolve()
        # Double-check: the resolved path must stay under skill_dir.
        try:
            target_file.relative_to(target.skill_dir.resolve())
        except ValueError:
            return f"Error: Path escapes the skill directory: '{file}'. Use a relative path within '{name}' (e.g. 'scripts/generate.js')."
    else:
        target_file = target.skill_file

    if not target_file.exists():
        if file is not None:
            return f"Error: File not found in skill '{name}': '{file}'. Check available paths with the SKILL.md documentation."
        return f"Error: SKILL.md not found for skill '{name}'."

    if target_file.is_dir():
        # List the directory contents to help discovery.
        try:
            entries = sorted(e.name for e in target_file.iterdir() if not e.name.startswith("."))
        except OSError as exc:
            return f"Error: Cannot list directory in '{name}': {exc}"
        if not entries:
            return f"Directory '{file}' in skill '{name}' is empty."
        return f"Directory listing for '{name}/{file}':\n" + "\n".join(f"  {e}" for e in entries)

    # Dispatch on file extension — plain-text, .docx, .pptx, or unsupported binary.
    content = _read_skill_file(target_file)
    if content.startswith("Error:"):
        return content

    if not content.strip():
        return f"Error: File '{file or 'SKILL.md'}' in skill '{name}' is empty."

    # Strip YAML frontmatter so the agent (and frontend) only see
    # the skill body content, not internal metadata fields.
    content = _FRONTMATTER_RE.sub("", content).lstrip("\n")

    # ── Auto-inject skill files into sandbox ────────────────────
    # When a skill contains binary files (.docx templates, images,
    # etc.), copy the entire skill directory into the sandbox so
    # the agent can execute scripts and use templates directly.
    if sandbox_provider is not None and runtime is not None and file is not None:
        injected: set[str] | None = runtime.state.get("injected_skills") if runtime.state is not None else None
        already_injected = injected is not None and name in injected
        if not already_injected and _skill_has_binary_files(target.skill_dir):
            content += _inject_skill_files(target.skill_dir, sandbox_provider, runtime)

    return content


def _decode_fetched_text(name: str, rel_path: str, data: bytes, *, strip_frontmatter: bool) -> str:
    """Decode a single in-memory skill file as UTF-8 text."""
    try:
        text = data.decode("utf-8")
    except UnicodeDecodeError:
        return f"Error: Cannot read '{rel_path}' in skill '{name}' as text — it appears to be a binary file.  The skill's files have been injected into the sandbox; use sandbox tools (read_file/bash) to access them."
    if strip_frontmatter:
        text = _FRONTMATTER_RE.sub("", text).lstrip("\n")
    if not text.strip():
        return f"Error: File '{rel_path}' in skill '{name}' is empty."
    return text


def _render_fetched_skill(name: str, files: list[tuple[str, bytes]], file: str | None) -> str:
    """Render a custom skill's ``SKILL.md`` or a supporting file from bytes.

    *files* is the ``(rel_path, bytes)`` list returned by the
    ``fetch_skill_files`` callback.  ``file is None`` returns ``SKILL.md``;
    otherwise the matching file (or a directory listing under a prefix) is
    returned.
    """
    by_rel: dict[str, bytes] = {rel: data for rel, data in files}

    if file is not None:
        clean = file.replace("\\", "/").strip("/")
        for segment in clean.split("/"):
            if segment in {"..", ""}:
                return f"Error: Path traversal is not allowed: '{file}'. Use a relative path within the skill directory (e.g. 'scripts/generate.js')."

        if clean in by_rel:
            return _decode_fetched_text(name, clean, by_rel[clean], strip_frontmatter=(clean == "SKILL.md"))

        # Directory listing: entries directly under "<clean>/".
        prefix = clean.rstrip("/") + "/"
        entries = sorted({rel[len(prefix) :].split("/")[0] for rel in by_rel if rel.startswith(prefix) and rel != prefix})
        if entries:
            return f"Directory listing for '{name}/{clean}/':\n" + "\n".join(f"  {e}" for e in entries)

        return f"Error: File not found in skill '{name}': '{file}'. Check available paths with the SKILL.md documentation."

    if "SKILL.md" in by_rel:
        return _decode_fetched_text(name, "SKILL.md", by_rel["SKILL.md"], strip_frontmatter=True)
    return f"Error: SKILL.md not found for skill '{name}'."


def _build_read_skill_tool(
    skills_path: Path | None,
    sandbox_provider: SandboxProvider | None = None,
    *,
    is_available: IsAvailableCallback | None = None,
    fetch_skill_files: FetchSkillFilesCallback | None = None,
):
    """Build and return a ``read_skill`` langchain tool.

    Built-in skills (bundled with the image) are read directly from the
    host filesystem at *skills_path*.  When the optional *is_available*
    and *fetch_skill_files* callbacks are supplied, non-built-in skill
    names are treated as **custom skills**:

    * *is_available(name)* decides whether the skill is in the current
      user's available set (own skills or added & approved marketplace
      skills) — if not, the read is rejected.
    * *fetch_skill_files(name)* downloads the skill's files from object
      storage, which are then injected into the sandbox at
      ``.skills/<name>/`` on first access (so the agent can execute
      scripts / reference templates) and read from memory.

    When the callbacks are omitted (SDK standalone), the tool falls back
    to the original behaviour: only filesystem skills are visible.

    Args:
        skills_path: Host filesystem path to the built-in skills root.
        sandbox_provider: Optional sandbox provider for injecting skill
            files into the agent's sandbox workspace.
        is_available: Optional async callback ``(name) -> bool``.
        fetch_skill_files: Optional async callback
            ``(name) -> list[tuple[rel_path, bytes]]``.
    """

    @tool("read_skill", parse_docstring=True)
    async def read_skill(
        name: str,
        file: str | None = None,
        runtime: _RuntimeType = None,
    ) -> str:
        """Load a skill's SKILL.md content or a supporting file by name.

        Use this tool before tackling complex tasks that match one
        of the available skills.  Built-in skills are listed in the
        system prompt; personal skills can be discovered with the
        ``list_skills`` tool.  Each SKILL.md contains best practices,
        workflows, frameworks, and references to additional resources.

        When the SKILL.md references other files (scripts, references,
        templates, etc.), call this tool again with the ``file``
        argument to read those files directly.

        Args:
            name: The skill name (e.g. "chart-visualization").
            file: Optional relative path to a file within the skill
                directory (e.g. "scripts/generate.js").  When omitted,
                the main SKILL.md is returned.
        """
        # ── Built-in skills: read from the filesystem ────────────────
        if skills_path is not None and skills_path.exists():
            builtin = load_skills(skills_path, enabled_only=True)
            target = next((s for s in builtin if s.name == name), None)
            if target is not None:
                return _read_builtin_skill(target, file, sandbox_provider, runtime)

        # ── Custom skills: resolve via the injected callbacks ─────────
        if is_available is None:
            # No callback → only built-in skills exist in this runtime.
            builtin = load_skills(skills_path, enabled_only=True) if skills_path is not None and skills_path.exists() else []
            available = ", ".join(s.name for s in builtin)
            return f"Error: Skill '{name}' not found. Available skills: {available or '(none)'}"

        try:
            allowed = await is_available(name)
        except Exception as exc:
            logger.warning("is_available(%r) raised: %s", name, exc)
            allowed = False

        if not allowed:
            return f"Error: Skill '{name}' is not available to the current user. It must be one of your own skills, or a marketplace skill you have added (and that an admin has approved)."

        if fetch_skill_files is None:
            return f"Error: Skill '{name}' is available, but no fetch_skill_files callback is configured."

        try:
            files = await fetch_skill_files(name)
        except Exception as exc:
            logger.warning("fetch_skill_files(%r) raised: %s", name, exc)
            return f"Error: Failed to load custom skill '{name}': {exc}"

        if not files:
            return f"Error: Skill '{name}' has no files in storage."

        # ── Inject into the sandbox (once per conversation) ─────────
        injection_msg = ""
        if sandbox_provider is not None and runtime is not None:
            injected = runtime.state.get("injected_skills") if runtime.state is not None else None
            if injected is None or name not in injected:
                injection_msg = _inject_skill_file_list(files, name, sandbox_provider, runtime)

        content = _render_fetched_skill(name, files, file)
        if content.startswith("Error:"):
            return content
        return content + injection_msg

    return read_skill


def _build_list_skills_tool(list_personal_skills: ListPersonalSkillsCallback):
    """Build and return a ``list_skills`` langchain tool.

    Returns the current user's personal skills (own + added & approved),
    fetched at runtime via the injected *list_personal_skills* callback.
    Built-in skills are intentionally excluded — they are already listed
    in the system prompt's ``<available_skills>`` block.
    """

    @tool("list_skills", parse_docstring=True)
    async def list_skills() -> str:
        """List the current user's personal skills.

        Built-in skills are already listed in the system prompt's
        ``<available_skills>`` block.  Use this tool to discover the
        user's **personal** skills (skills they created themselves, or
        marketplace skills they added), which are NOT in the system
        prompt.  After finding a relevant skill, load it with
        ``read_skill('<name>')``.
        """
        try:
            skills = await list_personal_skills()
        except Exception as exc:
            logger.warning("list_personal_skills() raised: %s", exc)
            return f"Error: Failed to list personal skills: {exc}"

        if not skills:
            return "当前没有个人技能（自己创建或已添加的技能）。"

        lines = ["当前用户可用的个人技能："]
        for name, description in skills:
            desc = description.strip().split("\n")[0] if description else ""
            suffix = f" — {desc}" if desc else ""
            lines.append(f"- {name}{suffix} — read with `read_skill('{name}')`")
        return "\n".join(lines)

    return list_skills


def make_skill_tools(
    *,
    skills_path: Path | None,
    sandbox_provider: SandboxProvider | None = None,
    is_available: IsAvailableCallback | None = None,
    fetch_skill_files: FetchSkillFilesCallback | None = None,
    list_personal_skills: ListPersonalSkillsCallback | None = None,
):
    """Create skill-related tools for the agent runtime.

    Args:
        skills_path: Host filesystem path to the built-in skills root.
        sandbox_provider: Optional sandbox provider.  When supplied,
            skill files (binary built-in files, or custom skills fetched
            from storage) are injected into the sandbox on first access.
        is_available: Optional async ``(name) -> bool`` callback deciding
            whether a custom skill is available to the current user.
        fetch_skill_files: Optional async callback downloading a custom
            skill's files from object storage.
        list_personal_skills: Optional async callback returning the
            current user's personal skills.  When supplied, a
            ``list_skills`` tool is also returned.

    Returns:
        A list of langchain :class:`BaseTool` instances: ``read_skill``,
        plus ``list_skills`` when *list_personal_skills* is supplied.
        Returns an empty list when *skills_path* does not exist and no
        personal-skills callback is supplied.
    """
    if (skills_path is None or not skills_path.exists()) and list_personal_skills is None:
        return []
    tools = [
        _build_read_skill_tool(
            skills_path,
            sandbox_provider,
            is_available=is_available,
            fetch_skill_files=fetch_skill_files,
        )
    ]
    if list_personal_skills is not None:
        tools.append(_build_list_skills_tool(list_personal_skills))
    return tools
